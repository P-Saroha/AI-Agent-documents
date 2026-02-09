"""
Document Ingestion Module
Handles loading and processing of various document types
"""
import os
from typing import List, Dict
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd


class DocumentProcessor:
    """Process various document types and extract text"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
    
    def load_pdf(self, file_path: str) -> List[Document]:
        """Extract text from PDF files"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            
            metadata = {
                "source": file_path,
                "type": "pdf",
                "pages": len(reader.pages)
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            raise Exception(f"Error loading PDF {file_path}: {str(e)}")
    
    def load_docx(self, file_path: str) -> List[Document]:
        """Extract text from DOCX files"""
        try:
            doc = DocxDocument(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            metadata = {
                "source": file_path,
                "type": "docx",
                "paragraphs": len(doc.paragraphs)
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            raise Exception(f"Error loading DOCX {file_path}: {str(e)}")
    
    def load_pptx(self, file_path: str) -> List[Document]:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            metadata = {
                "source": file_path,
                "type": "pptx",
                "slides": len(prs.slides)
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            raise Exception(f"Error loading PPTX {file_path}: {str(e)}")
    
    def load_excel(self, file_path: str) -> List[Document]:
        """Extract text from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            documents = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Process in batches of rows to avoid huge single document
                batch_size = 50  # 50 rows per chunk
                for start in range(0, len(df), batch_size):
                    batch = df.iloc[start:start + batch_size]
                    text = f"Sheet: {sheet_name} (Rows {start+1}-{start+len(batch)})\n"
                    text += batch.to_string()
                    
                    metadata = {
                        "source": file_path,
                        "type": "excel",
                        "sheet": sheet_name,
                        "rows": f"{start+1}-{start+len(batch)}"
                    }
                    documents.append(Document(page_content=text, metadata=metadata))
            
            return documents
        except Exception as e:
            raise Exception(f"Error loading Excel {file_path}: {str(e)}")
    
    def load_txt(self, file_path: str) -> List[Document]:
        """Extract text from text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            metadata = {
                "source": file_path,
                "type": "txt"
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            raise Exception(f"Error loading TXT {file_path}: {str(e)}")
    
    def process_document(self, file_path: str) -> List[Document]:
        """Process a document based on its extension"""
        file_path = str(file_path)
        extension = Path(file_path).suffix.lower()
        
        loaders = {
            '.pdf': self.load_pdf,
            '.docx': self.load_docx,
            '.pptx': self.load_pptx,
            '.xlsx': self.load_excel,
            '.xls': self.load_excel,
            '.txt': self.load_txt,
        }
        
        if extension not in loaders:
            raise ValueError(f"Unsupported file type: {extension}")
        
        # Load document
        documents = loaders[extension](file_path)
        
        # Split into chunks
        chunks = self.text_splitter.split_documents(documents)
        
        return chunks
    
    def process_multiple_documents(self, file_paths: List[str]) -> List[Document]:
        """Process multiple documents"""
        all_chunks = []
        
        for file_path in file_paths:
            try:
                chunks = self.process_document(file_path)
                all_chunks.extend(chunks)
            except Exception as e:
                pass
        
        return all_chunks
